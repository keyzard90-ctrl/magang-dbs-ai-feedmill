import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_text_from_shape(shape):
    text = ""
    if hasattr(shape, "text") and shape.text:
        text += shape.text + "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            text += get_text_from_shape(s)
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text_frame.text + " "
            text += "\n"
    return text

def extract_text(ppt_path):
    prs = Presentation(ppt_path)
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            print(get_text_from_shape(shape).strip())

if __name__ == "__main__":
    extract_text(sys.argv[1])
